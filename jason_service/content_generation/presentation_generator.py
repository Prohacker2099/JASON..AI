import json
import logging
import os
import re
import time
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
import random

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available - PowerPoint generation disabled")

@dataclass
class SlideContent:
    """Structured content for a presentation slide."""
    title: str
    content_points: List[str]
    speaker_notes: str
    image_query: str
    slide_type: str = "content"  # title, content, comparison, chart, quote
    transition: str = "fade"
    animation: str = "appear"

@dataclass
class PresentationStructure:
    """Complete presentation structure with metadata."""
    title: str
    subtitle: str
    author: str
    date: str
    slides: List[SlideContent]
    theme: str = "modern"
    color_scheme: str = "professional"
    total_slides: int = 0

class PromptExpander:
    """Advanced prompt expansion system for 10,000x content generation."""
    
    def __init__(self, local_llm=None):
        self.expansion_patterns = self._load_expansion_patterns()
        self.content_templates = self._load_content_templates()
        self.knowledge_domains = self._load_knowledge_domains()
        self.local_llm = local_llm

    def _load_expansion_patterns(self) -> Dict[str, List[str]]:
        """Load patterns for expanding prompts into comprehensive content."""
        return {
            'business': [
                "Market analysis and competitive landscape",
                "Financial projections and ROI calculations",
                "Risk assessment and mitigation strategies",
                "Implementation timeline and milestones",
                "Team structure and resource allocation",
                "Technology stack and infrastructure requirements",
                "Marketing and growth strategies",
                "Operational efficiency and process optimization",
                "Customer acquisition and retention",
                "Regulatory compliance and legal considerations"
            ],
            'technical': [
                "Architecture overview and system design",
                "Technical specifications and requirements",
                "Development methodology and best practices",
                "Testing strategies and quality assurance",
                "Security considerations and data protection",
                "Performance optimization and scalability",
                "Integration capabilities and APIs",
                "Maintenance and support procedures",
                "Documentation and knowledge transfer",
                "Future roadmap and technology trends"
            ],
            'educational': [
                "Learning objectives and outcomes",
                "Core concepts and theoretical foundations",
                "Practical applications and case studies",
                "Interactive elements and engagement strategies",
                "Assessment methods and evaluation criteria",
                "Resources and materials",
                "Prerequisites and target audience",
                "Learning path and progression",
                "Common challenges and solutions",
                "Further learning and advanced topics"
            ],
            'creative': [
                "Creative concept and vision",
                "Visual design and aesthetics",
                "Brand identity and messaging",
                "User experience and interface design",
                "Content strategy and storytelling",
                "Multi-platform adaptation",
                "Audience engagement and interaction",
                "Innovation and differentiation",
                "Production workflow and timeline",
                "Measuring success and impact"
            ]
        }
    
    def _load_content_templates(self) -> Dict[str, Dict]:
        """Load templates for different types of content."""
        return {
            'introduction': {
                'word_count': (80, 110),
                'structure': ['hook', 'context', 'thesis', 'preview'],
                'elements': ['attention_grabber', 'background', 'purpose', 'roadmap']
            },
            'content': {
                'word_count': (80, 110),
                'structure': ['main_point', 'supporting_evidence', 'examples', 'implications'],
                'elements': ['key_concept', 'data_points', 'case_study', 'takeaways']
            },
            'conclusion': {
                'word_count': (80, 110),
                'structure': ['summary', 'key_insights', 'call_to_action', 'future_outlook'],
                'elements': ['recap', 'insights', 'action_items', 'vision']
            }
        }
    
    def _load_knowledge_domains(self) -> Dict[str, List[str]]:
        """Load domain-specific knowledge for content enrichment."""
        return {
            'technology': [
                "Artificial Intelligence and Machine Learning",
                "Cloud Computing and Infrastructure",
                "Cybersecurity and Data Protection",
                "Blockchain and Distributed Systems",
                "Internet of Things (IoT)",
                "Quantum Computing",
                "5G and Next-Gen Networks",
                "Edge Computing",
                "DevOps and Automation",
                "Data Science and Analytics"
            ],
            'business': [
                "Strategic Planning and Management",
                "Digital Transformation",
                "Customer Experience Management",
                "Supply Chain Optimization",
                "Financial Technology (FinTech)",
                "E-commerce and Retail",
                "Healthcare and Life Sciences",
                "Manufacturing and Industry 4.0",
                "Real Estate and Construction",
                "Energy and Sustainability"
            ],
            'education': [
                "Online Learning Platforms",
                "Educational Technology",
                "Curriculum Development",
                "Assessment and Evaluation",
                "Learning Analytics",
                "Virtual and Augmented Reality in Education",
                "Personalized Learning",
                "STEM Education",
                "Professional Development",
                "Educational Policy and Administration"
            ]
        }

    def expand_prompt_10000x(self, original_prompt: str, target_slides: int = 20) -> PresentationStructure:
        """Expand a simple prompt into comprehensive 10,000x content using LLM if available."""
        logging.info(f"Expanding prompt: '{original_prompt}' into {target_slides} slides")
        
        # 1. Try LLM Generation First
        if self.local_llm and self.local_llm.available:
            try:
                logging.info("Attempting LLM-based structure generation...")
                structure = self._generate_structure_via_llm(original_prompt, target_slides)
                if structure:
                     logging.info(f"LLM generated {len(structure.slides)} slides.")
                     return structure
            except Exception as e:
                logging.error(f"LLM generation failed, falling back to templates: {e}")

        # 2. Fallback to Template Logic
        logging.info("Using template-based generation.")
        
        # Analyze prompt to determine domain and type
        domain = self._classify_domain(original_prompt)
        content_type = self._classify_content_type(original_prompt)
        
        # Generate comprehensive structure
        structure = self._generate_structure(original_prompt, domain, content_type, target_slides)
        
        # Expand each slide with detailed content
        for i, slide in enumerate(structure.slides):
            slide.content_points = self._expand_slide_content(slide, domain, content_type)
            slide.speaker_notes = self._generate_speaker_notes(slide, domain)
            slide.image_query = self._generate_image_query(slide, domain)
        
        structure.total_slides = len(structure.slides)
        
        logging.info(f"Generated {structure.total_slides} slides with comprehensive content")
        return structure

    def _generate_structure_via_llm(self, prompt: str, target_slides: int) -> Optional[PresentationStructure]:
        """Generate full presentation structure using Local LLM."""
        
        schema = """
        {
            "title": "string",
            "subtitle": "string",
            "slides": [
                {
                    "title": "string",
                    "content_points": ["string", "string", "string"],
                    "speaker_notes": "string",
                    "image_query": "string (for search)",
                    "slide_type": "title|content|conclusion"
                }
            ]
        }
        """
        
        llm_prompt = f"""
        Act as an expert presentation designer. Create a comprehensive {target_slides}-slide presentation outline for the topic: "{prompt}".
        Expand this into a professional, deep-dive presentation.
        Include detailed content points, speaker notes, and image search queries for each slide.
        Output ONLY valid JSON matching this schema:
        {schema}
        """
        
        data = self.local_llm.generate_json(llm_prompt, schema)
        if not data:
            return None
            
        slides = []
        for s in data.get('slides', []):
            slides.append(SlideContent(
                title=s.get('title', 'Untitled'),
                content_points=s.get('content_points', []),
                speaker_notes=s.get('speaker_notes', ''),
                image_query=s.get('image_query', ''),
                slide_type=s.get('slide_type', 'content')
            ))
            
        return PresentationStructure(
            title=data.get('title', prompt),
            subtitle=data.get('subtitle', 'Generated by JASON AI'),
            author="JASON AI",
            date=datetime.now().strftime("%B %d, %Y"),
            slides=slides,
            total_slides=len(slides)
        )
    
    def _classify_domain(self, prompt: str) -> str:
        """Classify the domain of the prompt."""
        prompt_lower = prompt.lower()
        
        domain_keywords = {
            'technology': ['ai', 'ml', 'software', 'tech', 'computer', 'digital', 'data', 'system'],
            'business': ['business', 'market', 'finance', 'strategy', 'management', 'company', 'revenue'],
            'education': ['education', 'learning', 'teaching', 'student', 'academic', 'curriculum', 'school']
        }
        
        for domain, keywords in domain_keywords.items():
            if any(keyword in prompt_lower for keyword in keywords):
                return domain
        
        return 'general'
    
    def _classify_content_type(self, prompt: str) -> str:
        """Classify the type of content needed."""
        prompt_lower = prompt.lower()
        
        if any(word in prompt_lower for word in ['presentation', 'slides', 'talk']):
            return 'presentation'
        elif any(word in prompt_lower for word in ['report', 'analysis', 'study']):
            return 'report'
        elif any(word in prompt_lower for word in ['tutorial', 'guide', 'how to']):
            return 'tutorial'
        else:
            return 'presentation'
    
    def _generate_structure(self, prompt: str, domain: str, content_type: str, target_slides: int) -> PresentationStructure:
        """Generate the overall presentation structure."""
        
        # Create title and metadata
        structure = PresentationStructure(
            title=prompt.title(),
            subtitle=f"Comprehensive Analysis and Strategic Insights",
            author="JASON - AI Architect",
            date=datetime.now().strftime("%B %d, %Y"),
            slides=[],
            theme="modern",
            color_scheme="professional"
        )
        
        # Generate slide sequence based on content type
        if content_type == 'presentation':
            slides_config = self._generate_presentation_slides(prompt, domain, target_slides)
        elif content_type == 'report':
            slides_config = self._generate_report_slides(prompt, domain, target_slides)
        else:  # tutorial
            slides_config = self._generate_tutorial_slides(prompt, domain, target_slides)
        
        # Create slide objects
        for slide_config in slides_config:
            slide = SlideContent(
                title=slide_config['title'],
                content_points=slide_config.get('content_points', []),
                speaker_notes="",
                image_query=slide_config.get('image_query', ''),
                slide_type=slide_config.get('slide_type', 'content'),
                transition=slide_config.get('transition', 'fade'),
                animation=slide_config.get('animation', 'appear')
            )
            structure.slides.append(slide)
        
        return structure
    
    def _generate_presentation_slides(self, prompt: str, domain: str, target_slides: int) -> List[Dict]:
        """Generate slides for presentation type."""
        slides = []
        
        # Title slide
        slides.append({
            'title': prompt.title(),
            'slide_type': 'title',
            'image_query': f"professional {domain} concept header image"
        })
        
        # Overview slide
        slides.append({
            'title': "Executive Overview",
            'slide_type': 'content',
            'image_query': f"{domain} strategic overview diagram"
        })
        
        # Domain-specific expansion patterns
        patterns = self.expansion_patterns.get(domain, self.expansion_patterns['business'])
        num_content_slides = target_slides - 3  # Account for title, overview, conclusion
        
        # Generate content slides
        for i in range(min(num_content_slides, len(patterns))):
            pattern = patterns[i % len(patterns)]
            slides.append({
                'title': pattern,
                'slide_type': 'content',
                'image_query': f"{domain} {pattern.lower()} visualization"
            })
        
        # Add more slides if needed
        while len(slides) < target_slides - 1:
            additional_topics = [
                "Key Success Factors",
                "Implementation Challenges",
                "Best Practices and Recommendations",
                "Future Trends and Opportunities",
                "Measuring Success and ROI"
            ]
            topic = additional_topics[(len(slides) - 2) % len(additional_topics)]
            slides.append({
                'title': topic,
                'slide_type': 'content',
                'image_query': f"{domain} {topic.lower()} infographic"
            })
        
        # Conclusion slide
        slides.append({
            'title': "Conclusion and Next Steps",
            'slide_type': 'conclusion',
            'image_query': f"{domain} future vision concept"
        })
        
        return slides[:target_slides]
    
    def _generate_report_slides(self, prompt: str, domain: str, target_slides: int) -> List[Dict]:
        """Generate slides for report type."""
        slides = []
        
        # Title and methodology
        slides.extend([
            {'title': prompt.title(), 'slide_type': 'title'},
            {'title': "Research Methodology", 'slide_type': 'content'},
            {'title': "Key Findings", 'slide_type': 'content'},
        ])
        
        # Data and analysis slides
        for i in range(target_slides - 5):
            slides.append({
                'title': f"Analysis Insight {i+1}",
                'slide_type': 'content',
                'image_query': f"{domain} data visualization chart"
            })
        
        # Conclusion
        slides.extend([
            {'title': "Implications and Impact", 'slide_type': 'content'},
            {'title': "Recommendations", 'slide_type': 'conclusion'}
        ])
        
        return slides
    
    def _generate_tutorial_slides(self, prompt: str, domain: str, target_slides: int) -> List[Dict]:
        """Generate slides for tutorial type."""
        slides = []
        
        # Introduction
        slides.extend([
            {'title': prompt.title(), 'slide_type': 'title'},
            {'title': "Learning Objectives", 'slide_type': 'content'},
            {'title': "Prerequisites and Setup", 'slide_type': 'content'},
        ])
        
        # Tutorial steps
        for i in range(target_slides - 4):
            slides.append({
                'title': f"Step {i+1}: Key Concept",
                'slide_type': 'content',
                'image_query': f"{domain} tutorial step {i+1} diagram"
            })
        
        # Conclusion
        slides.append({'title': "Summary and Resources", 'slide_type': 'conclusion'})
        
        return slides
    
    def _expand_slide_content(self, slide: SlideContent, domain: str, content_type: str) -> List[str]:
        """Expand slide content with detailed points."""
        template = self.content_templates.get(slide.slide_type, self.content_templates['content'])
        min_words, max_words = template['word_count']
        
        # Generate content based on slide type and domain
        if slide.slide_type == 'title':
            return [f"Comprehensive overview of {slide.title.lower()}"]
        
        # Generate 3-5 detailed bullet points
        content_points = []
        num_points = random.randint(3, 5)
        
        for i in range(num_points):
            point = self._generate_content_point(slide.title, domain, i, num_points)
            content_points.append(point)
        
        return content_points
    
    def _generate_content_point(self, slide_title: str, domain: str, index: int, total_points: int) -> str:
        """Generate a single content point."""
        # Context-aware content generation
        if index == 0:
            # First point - introduce the concept
            return f"Understanding the fundamental principles and core concepts of {slide_title.lower()} within the {domain} ecosystem"
        elif index == total_points - 1:
            # Last point - implications or next steps
            return f"Exploring the strategic implications and future opportunities related to {slide_title.lower()}"
        else:
            # Middle points - detailed aspects
            aspects = [
                "Critical analysis and evaluation",
                "Practical implementation strategies",
                "Best practices and methodologies",
                "Performance metrics and success indicators",
                "Risk assessment and mitigation approaches"
            ]
            aspect = aspects[index % len(aspects)]
            return f"{aspect} for {slide_title.lower()} with measurable outcomes and actionable insights"
    
    def _generate_speaker_notes(self, slide: SlideContent, domain: str) -> str:
        """Generate detailed speaker notes."""
        notes = f"Speaker Notes for {slide.title}:\n\n"
        notes += f"This slide focuses on {slide.title.lower()} within the {domain} context. "
        notes += f"Key talking points include:\n"
        
        for i, point in enumerate(slide.content_points, 1):
            notes += f"{i}. {point}\n"
        
        notes += f"\nAnticipated questions and discussion points:\n"
        notes += f"- How does this concept apply to real-world scenarios?\n"
        notes += f"- What are the implementation challenges?\n"
        notes += f"- How can we measure success in this area?\n"
        
        return notes
    
    def _generate_image_query(self, slide: SlideContent, domain: str) -> str:
        """Generate image search query for the slide."""
        base_query = f"professional {domain} {slide.title.lower()}"
        
        # Add context based on slide type
        if slide.slide_type == 'title':
            return f"{base_query} header banner modern design"
        elif slide.slide_type == 'conclusion':
            return f"{base_query} future vision concept"
        else:
            return f"{base_query} infographic diagram visualization"

class PresentationGenerator:
    """Production-ready PowerPoint generator with 10,000x prompt expansion."""
    
    def __init__(self, uspt_model=None, local_llm=None, image_search_func=None, output_dir="/tmp"):
        self.uspt_model = uspt_model
        self.local_llm = local_llm
        self.image_search_func = image_search_func
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        
        self.prompt_expander = PromptExpander(local_llm=self.local_llm)
        
        if not PPTX_AVAILABLE:
            logging.error("python-pptx not available - PowerPoint generation disabled")
        
        logging.info(f"Presentation Generator initialized with output directory: {output_dir}")
    
    def create_presentation(self, prompt: str, filename: Optional[str] = None, target_slides: int = 20) -> str:
        """Create a comprehensive PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PowerPoint generation")
        
        logging.info(f"Creating presentation: '{prompt}' with {target_slides} slides")
        
        # Expand prompt into comprehensive structure
        structure = self.prompt_expander.expand_prompt_10000x(prompt, target_slides)
        
        # Generate filename if not provided
        if not filename:
            safe_title = re.sub(r'[^\w\s-]', '', structure.title).strip()
            safe_title = re.sub(r'[-\s]+', '_', safe_title)
            filename = f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        output_path = self.output_dir / filename
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Add slides
        for i, slide in enumerate(structure.slides):
            self._add_slide_to_presentation(prs, slide, i == 0)
        
        # Save presentation
        prs.save(str(output_path))
        
        # Save metadata
        self._save_presentation_metadata(structure, output_path.with_suffix('.json'))
        
        logging.info(f"Presentation saved to: {output_path}")
        return str(output_path)
    
    def _add_slide_to_presentation(self, prs: Presentation, slide: SlideContent, is_title_slide: bool = False):
        """Add a slide to the PowerPoint presentation."""
        if is_title_slide:
            # Title slide
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            ppt_slide = prs.slides.add_slide(slide_layout)
            ppt_slide.shapes.title.text = slide.title
            ppt_slide.placeholders[1].text = structure.subtitle
        else:
            # Content slide
            slide_layout = prs.slide_layouts[1]  # Title and content
            ppt_slide = prs.slides.add_slide(slide_layout)
            ppt_slide.shapes.title.text = slide.title
            
            # Add content points
            if slide.content_points:
                content_placeholder = ppt_slide.placeholders[1]
                tf = content_placeholder.text_frame
                tf.clear()
                
                for point in slide.content_points:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.name = 'Calibri'
        
        # Add speaker notes
        if slide.speaker_notes:
            notes_slide = ppt_slide.notes_slide
            notes_slide.notes_text_frame.text = slide.speaker_notes
        
        # Add image if available
        if slide.image_query and self.image_search_func:
            try:
                image_path = self.image_search_func(slide.image_query)
                if image_path and os.path.exists(image_path):
                    # Add image to slide (position and size will vary by layout)
                    left = Inches(6.0)
                    top = Inches(2.0)
                    width = Inches(3.0)
                    height = Inches(2.0)
                    ppt_slide.shapes.add_picture(image_path, left, top, width, height)
            except Exception as e:
                logging.warning(f"Failed to add image for query '{slide.image_query}': {e}")
    
    def _save_presentation_metadata(self, structure: PresentationStructure, metadata_path: Path):
        """Save presentation metadata to JSON file."""
        metadata = {
            'title': structure.title,
            'subtitle': structure.subtitle,
            'author': structure.author,
            'date': structure.date,
            'theme': structure.theme,
            'color_scheme': structure.color_scheme,
            'total_slides': structure.total_slides,
            'slides': [
                {
                    'title': slide.title,
                    'slide_type': slide.slide_type,
                    'content_points': slide.content_points,
                    'image_query': slide.image_query,
                    'transition': slide.transition,
                    'animation': slide.animation
                }
                for slide in structure.slides
            ],
            'generated_at': datetime.now().isoformat()
        }
        
        with open(metadata_path, 'w') as f:
            json.dump(metadata, f, indent=2)
        
        logging.info(f"Presentation metadata saved to: {metadata_path}")
    
    def get_presentation_templates(self) -> List[str]:
        """Get available presentation templates."""
        return [
            "modern_business",
            "academic_research", 
            "technical_documentation",
            "creative_portfolio",
            "executive_summary",
            "training_material"
        ]
    
    def batch_generate_presentations(self, prompts: List[str], target_slides: int = 20) -> List[str]:
        """Generate multiple presentations from a list of prompts."""
        results = []
        
        for prompt in prompts:
            try:
                output_path = self.create_presentation(prompt, target_slides=target_slides)
                results.append(output_path)
                logging.info(f"Generated presentation for: {prompt}")
            except Exception as e:
                logging.error(f"Failed to generate presentation for '{prompt}': {e}")
                results.append(None)
        
        return results

# Example usage
if __name__ == '__main__':
    def sample_image_search(query):
        """Sample image search function."""
        return None  # Placeholder
    
    generator = PresentationGenerator(
        uspt_model=None,
        image_search_func=sample_image_search,
        output_dir="/tmp"
    )
    
    # Test presentation generation
    prompt = "Artificial Intelligence in Modern Business"
    output_path = generator.create_presentation(prompt)
    print(f"Presentation generated: {output_path}")
